"""多格式文档解析器 — 支持 PDF / Word / Excel / PPT / HTML / CSV / Markdown / 纯文本。

设计原则：
- 所有解析器可选依赖：未安装对应库时给出明确安装提示，不会 import 报错
- 统一返回 (text: str, metadata: dict) 元组
- 自动检测编码（UTF-8 / GBK / GB2312）
"""

from __future__ import annotations

import csv
from pathlib import Path
from typing import Any

# ------------------------------------------------------------------
# 统一入口
# ------------------------------------------------------------------


def parse_file(file_path: str | Path) -> tuple[str, dict[str, Any]]:
    """自动检测文件类型并解析。

    Returns:
        (提取的全文文本, 元数据字典)
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    if not file_path.is_file():
        raise ValueError(f"不是文件: {file_path}")

    suffix = file_path.suffix.lower()
    size_bytes = file_path.stat().st_size

    if suffix in (".pdf",):
        text = parse_pdf(file_path)
    elif suffix in (".docx", ".doc"):
        text = parse_docx(file_path)
    elif suffix in (".pptx", ".ppt"):
        text = parse_pptx(file_path)
    elif suffix in (".xlsx", ".xls", ".xlsm"):
        text = parse_excel(file_path)
    elif suffix in (".html", ".htm"):
        text = parse_html(file_path)
    elif suffix in (".csv", ".tsv"):
        text = parse_csv(file_path)
    elif suffix in (
        ".md",
        ".markdown",
        ".txt",
        ".log",
        ".json",
        ".xml",
        ".yaml",
        ".yml",
        ".py",
        ".js",
        ".ts",
        ".java",
        ".go",
        ".rs",
        ".cpp",
        ".c",
        ".h",
    ):
        text = parse_text(file_path)
    else:
        # 未知类型 fallback 为纯文本
        text = parse_text(file_path)

    return text, {
        "source": file_path.name,
        "path": str(file_path.resolve()),
        "type": suffix.lstrip("."),
        "size_bytes": size_bytes,
    }


# ------------------------------------------------------------------
# 解析器实现
# ------------------------------------------------------------------


def parse_pdf(file_path: str | Path) -> str:
    """解析 PDF 文件，提取文本内容。

    优先使用 pdfplumber（表格识别更好），
    回退到 PyPDF2。
    """
    text_parts: list[str] = []
    file_path = Path(file_path)

    # 尝试 pdfplumber
    try:
        import pdfplumber

        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
                # 也尝试提取表格
                for table in page.extract_tables():
                    if table:
                        rows = [" | ".join(str(cell or "") for cell in row) for row in table]
                        text_parts.append("\n".join(rows))
        return "\n\n".join(text_parts)
    except ImportError:
        pass
    except Exception:
        pass

    # 回退到 PyPDF2
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(file_path))
        for pdf_page in reader.pages:
            page_text = pdf_page.extract_text()
            if page_text:
                text_parts.append(page_text)
        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError(
            "无法解析 PDF 文件。请安装 pdfplumber 或 PyPDF2:\n  pip install pdfplumber\n  pip install PyPDF2"
        ) from None


def parse_docx(file_path: str | Path) -> str:
    """解析 Word (.docx) 文件。"""
    try:
        import docx
    except ImportError:
        raise ImportError("无法解析 Word 文件。请安装 python-docx:\n  pip install python-docx") from None

    doc = docx.Document(str(file_path))
    paragraphs: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            paragraphs.append(" | ".join(cells))

    return "\n\n".join(paragraphs)


def parse_pptx(file_path: str | Path) -> str:
    """解析 PowerPoint (.pptx) 文件。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("无法解析 PPT 文件。请安装 python-pptx:\n  pip install python-pptx") from None

    prs = Presentation(str(file_path))
    slides_text: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"--- 幻灯片 {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_parts.append(para.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    slide_parts.append(" | ".join(cells))
        slides_text.append("\n".join(slide_parts))

    return "\n\n".join(slides_text)


def parse_excel(file_path: str | Path) -> str:
    """解析 Excel 文件，每个 sheet 一个段落。"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("无法解析 Excel 文件。请安装 openpyxl:\n  pip install openpyxl") from None

    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    sheets_text: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = [f"--- Sheet: {sheet_name} ---"]

        for row in ws.iter_rows(values_only=True):
            row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_str.strip(" |"):
                rows.append(row_str)

        sheets_text.append("\n".join(rows))

    wb.close()
    return "\n\n".join(sheets_text)


def parse_html(file_path: str | Path) -> str:
    """解析 HTML 文件，提取纯文本。"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("无法解析 HTML 文件。请安装 beautifulsoup4:\n  pip install beautifulsoup4") from None

    with open(file_path, encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    # 移除 script/style 标签
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    return soup.get_text(separator="\n", strip=True)


def parse_csv(file_path: str | Path) -> str:
    """解析 CSV / TSV 文件。"""
    file_path = Path(file_path)
    delimiter = "\t" if file_path.suffix.lower() == ".tsv" else ","

    with open(file_path, encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]

    return "\n".join(rows)


def parse_text(file_path: str | Path) -> str:
    """解析纯文本文件，自动检测编码。"""
    file_path = Path(file_path)

    # 尝试 UTF-8
    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        pass

    # 尝试 GBK / GB2312
    for enc in ("gbk", "gb2312", "latin-1"):
        try:
            return file_path.read_text(encoding=enc, errors="replace")
        except Exception:
            continue

    # 终极 fallback
    return file_path.read_text(encoding="utf-8", errors="replace")
